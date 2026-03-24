"""Document and media parsing utilities used before retrieval and prompting."""

from __future__ import annotations

import csv
from dataclasses import dataclass
import hashlib
import json
from io import StringIO
from pathlib import Path

from .config import DOCUMENT_EXTENSIONS, IMAGE_EXTENSIONS, TEXT_EXTENSIONS, VIDEO_EXTENSIONS


def detect_media_kind(path: Path) -> str:
    """Classify a file path into the coarse media types used by the app."""

    suffix = path.suffix.lower()
    if suffix in IMAGE_EXTENSIONS:
        return "image"
    if suffix in VIDEO_EXTENSIONS:
        return "video"
    if suffix in TEXT_EXTENSIONS or suffix in DOCUMENT_EXTENSIONS:
        return "document"
    return "binary"


@dataclass(slots=True)
class ParsedDocument:
    """Text and derived visuals extracted from one document."""

    text: str
    image_paths: list[Path]
    parser_name: str


def extract_text_from_path(path: Path) -> str:
    """Read a supported document into plain text for retrieval."""

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _read_pdf(path)
    if suffix == ".docx":
        return _read_docx(path)
    if suffix == ".pptx":
        return _read_pptx(path)
    if suffix == ".xlsx":
        return _read_xlsx(path)
    if suffix == ".ipynb":
        return _read_notebook(path)
    if suffix == ".csv":
        return _read_csv(path)
    return path.read_text(encoding="utf-8", errors="ignore")


def parse_document_for_chat(
    path: Path,
    *,
    cache_root: Path,
    max_images: int = 8,
    max_rendered_pages: int = 2,
) -> ParsedDocument:
    """Prepare a document for chat by extracting text and optional visuals."""

    suffix = path.suffix.lower()
    if suffix != ".pdf":
        return ParsedDocument(
            text=extract_text_from_path(path),
            image_paths=[],
            parser_name="python-reader",
        )
    return _parse_pdf_for_chat(
        path,
        cache_root=cache_root,
        max_images=max_images,
        max_rendered_pages=max_rendered_pages,
    )


def sample_video_frames_for_chat(
    path: Path,
    *,
    cache_root: Path,
    max_frames: int = 8,
) -> list[Path]:
    """Extract a small cached set of evenly spaced frames from a video file."""

    if max_frames <= 0:
        return []

    cache_dir = cache_root / _video_cache_key(path, max_frames=max_frames)
    cache_dir.mkdir(parents=True, exist_ok=True)
    manifest_path = cache_dir / "manifest.json"

    if manifest_path.exists():
        try:
            payload = json.loads(manifest_path.read_text(encoding="utf-8"))
            frame_paths = [
                cache_dir / name
                for name in payload.get("frame_files", [])
                if (cache_dir / name).exists()
            ]
            if frame_paths:
                return frame_paths
        except Exception:
            pass

    from decord import VideoReader, cpu
    import numpy as np
    from PIL import Image

    reader = VideoReader(str(path), ctx=cpu(0))
    frame_count = len(reader)
    if frame_count <= 0:
        return []

    indices = np.linspace(0, frame_count - 1, num=min(max_frames, frame_count), dtype=int).tolist()
    batch = reader.get_batch(indices).asnumpy()

    frame_paths: list[Path] = []
    for output_index, (frame_index, frame_array) in enumerate(zip(indices, batch), start=1):
        target_path = cache_dir / f"frame-{output_index:02d}-at-{frame_index:06d}.png"
        Image.fromarray(frame_array).save(target_path, format="PNG")
        frame_paths.append(target_path)

    manifest = {
        "source_path": str(path.resolve()),
        "parser_name": "decord",
        "frame_files": [item.name for item in frame_paths],
    }
    manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    return frame_paths


def _read_pdf(path: Path) -> str:
    """Extract text from a PDF without generating chat visuals."""

    import fitz

    document = fitz.open(str(path))
    try:
        pages: list[str] = []
        for page_index in range(document.page_count):
            text = document.load_page(page_index).get_text("text").strip()
            if text:
                pages.append(f"Page {page_index + 1}\n{text}")
        return "\n\n".join(pages)
    finally:
        document.close()


def _parse_pdf_for_chat(
    path: Path,
    *,
    cache_root: Path,
    max_images: int,
    max_rendered_pages: int,
) -> ParsedDocument:
    """Parse a PDF into text plus extracted or rendered page images."""

    cache_dir = cache_root / _document_cache_key(
        path,
        max_images=max_images,
        max_rendered_pages=max_rendered_pages,
    )
    cache_dir.mkdir(parents=True, exist_ok=True)
    manifest_path = cache_dir / "manifest.json"
    text_path = cache_dir / "text.txt"

    if manifest_path.exists() and text_path.exists():
        try:
            payload = json.loads(manifest_path.read_text(encoding="utf-8"))
            image_paths = [
                cache_dir / name
                for name in payload.get("image_files", [])
                if (cache_dir / name).exists()
            ]
            return ParsedDocument(
                text=text_path.read_text(encoding="utf-8", errors="ignore"),
                image_paths=image_paths,
                parser_name=str(payload.get("parser_name", "pymupdf")),
            )
        except Exception:
            pass

    import fitz

    text_sections: list[str] = []
    image_paths: list[Path] = []
    seen_hashes: set[str] = set()

    document = fitz.open(str(path))
    try:
        for page_index in range(document.page_count):
            page = document.load_page(page_index)
            text = page.get_text("text").strip()
            if text:
                text_sections.append(f"Page {page_index + 1}\n{text}")

            # Prefer true embedded images so multimodal models see the original figures when possible.
            for image_index, image_ref in enumerate(page.get_images(full=True), start=1):
                if len(image_paths) >= max_images:
                    break
                extracted = document.extract_image(image_ref[0])
                image_bytes = extracted.get("image")
                if not image_bytes:
                    continue
                digest = hashlib.sha1(image_bytes).hexdigest()
                if digest in seen_hashes:
                    continue
                seen_hashes.add(digest)
                extension = str(extracted.get("ext") or "png").lower()
                target_path = cache_dir / f"page-{page_index + 1:03d}-image-{image_index:02d}.{extension}"
                target_path.write_bytes(image_bytes)
                image_paths.append(target_path)

        # If the PDF has no extractable figures, render a few page previews instead of giving the model no visuals.
        if not image_paths and max_rendered_pages > 0:
            for page_index in range(min(document.page_count, max_rendered_pages)):
                page = document.load_page(page_index)
                pixmap = page.get_pixmap(dpi=160, alpha=False)
                target_path = cache_dir / f"page-{page_index + 1:03d}-preview.png"
                pixmap.save(str(target_path))
                image_paths.append(target_path)
    finally:
        document.close()

    text = "\n\n".join(text_sections).strip()
    text_path.write_text(text, encoding="utf-8")
    manifest = {
        "source_path": str(path.resolve()),
        "parser_name": "pymupdf",
        "image_files": [item.name for item in image_paths],
    }
    manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    return ParsedDocument(text=text, image_paths=image_paths, parser_name="pymupdf")


def _document_cache_key(path: Path, *, max_images: int, max_rendered_pages: int) -> str:
    """Fingerprint a parsed-document cache entry by file identity and parser budget."""

    stat = path.stat()
    fingerprint = f"{path.resolve()}:{stat.st_mtime_ns}:{stat.st_size}:{max_images}:{max_rendered_pages}"
    return hashlib.sha1(fingerprint.encode("utf-8")).hexdigest()


def _video_cache_key(path: Path, *, max_frames: int) -> str:
    """Fingerprint a sampled-video cache entry by file identity and frame budget."""

    stat = path.stat()
    fingerprint = f"{path.resolve()}:{stat.st_mtime_ns}:{stat.st_size}:{max_frames}"
    return hashlib.sha1(fingerprint.encode("utf-8")).hexdigest()


def _read_docx(path: Path) -> str:
    """Extract paragraph text from a DOCX file."""

    from docx import Document

    document = Document(str(path))
    return "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())


def _read_pptx(path: Path) -> str:
    """Flatten slide text from a PowerPoint deck into plain text."""

    from pptx import Presentation

    presentation = Presentation(str(path))
    chunks: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        chunks.append(f"Slide {slide_index}")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                chunks.append(text)
    return "\n".join(chunks)


def _read_xlsx(path: Path) -> str:
    """Render worksheet cell values into a line-oriented text view."""

    from openpyxl import load_workbook

    workbook = load_workbook(path, data_only=True, read_only=True)
    lines: list[str] = []
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        lines.append(f"Sheet: {sheet_name}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell) for cell in row if cell not in (None, "")]
            if values:
                lines.append(" | ".join(values))
    return "\n".join(lines)


def _read_notebook(path: Path) -> str:
    """Extract cell contents from a Jupyter notebook JSON file."""

    payload = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
    lines: list[str] = []
    for cell in payload.get("cells", []):
        cell_type = cell.get("cell_type", "unknown")
        source = "".join(cell.get("source", []))
        if source.strip():
            lines.append(f"[{cell_type}]\n{source}")
    return "\n\n".join(lines)


def _read_csv(path: Path) -> str:
    """Round-trip a CSV file into normalized text for retrieval."""

    with path.open("r", encoding="utf-8", errors="ignore", newline="") as handle:
        reader = csv.reader(handle)
        buffer = StringIO()
        writer = csv.writer(buffer)
        for row in reader:
            writer.writerow(row)
        return buffer.getvalue()
